from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.util import Inches
from jira import JIRA
from tkinter import *
from tkinter import messagebox
from tkinter import simpledialog
from PIL import Image
from ConfReports import ConfReports
from Connect import Connect


import configparser
import copy
import qrcode
import six
import re
import os

global t_var
global root

# Open Parameters
def Parameters():
    os.startfile(CurrentDir + '\ini\Parameters.ini')
    return

def Openppt():
    os.startfile(CurrentDir + config['PPTX'].get('Output'))
    return

# ********************************
# Functions to call confluence Reports
# ********************************
def PIBoard():
    # Init text variable for lower Frame
    Status_var.set("Generating PI Board Report")
    t_var.set("Confluence Report in Progress")
    root.update()
    t_var.set(ConfReports.PIBoard())
    root.update()
    return

def TeamBoard():
    # Init text variable for lower Frame
    Status_var.set("Generating Team Reports")
    t_var.set("Confluence Report in Progress")
    root.update()
    ConfReports.TeamBoard()
    t_var.set("Confluence Report Completed")
    root.update()
    return

def AssigneeBoard():
    # Init text variable for lower Frame
    Status_var.set("Generating Assignee Reports")
    t_var.set("Confluence Report in Progress")
    root.update()
    ConfReports.AssigneeBoard()
    t_var.set("Confluence Report Completed")
    root.update()
    return

def PendingFeatures():
    # Init text variable for lower Frame
    Status_var.set("Generating PI Pending Scope")
    t_var.set("Confluence Report in Progress")
    root.update()
    ConfReports.PendingFeatures()
    t_var.set("Confluence Report Completed")
    root.update()
    return


# ********************************
# Menu Functions
# ********************************

def CheckJIRA():
    # Read Parameters to connect to JIRA
    jira_server = config['CONNECT'].get('JIRAserver')
    jira_user = config['CONNECT'].get('Username')
    jira_password = Connect.decryptPass()

    # Update Screen Labels
    Status_var.set("Connecting JIRA")
    root.update()

    try:
        # Connect to JIRA
        jira_server = {'server': jira_server}
        jira = JIRA(options=jira_server, basic_auth=(jira_user, jira_password))
        messagebox.showinfo(title="Information", message="Connected Successfully")
        Status_var.set("Click on Start Button")
        root.update_idletasks()
    except:
        messagebox.showinfo(title="Error", message="Connection Error")
        # Update Screen Labels
        Status_var.set("ERROR Connecting JIRA")
        final_str = 'Check Server: %s\n' \
                    'Check Username: %s \nCheck Password: %s' \
                    '\nChange these values in the Config Menu' % \
                    (config['CONNECT'].get('JIRAserver'), config['CONNECT'].get('Username'), Connect.decryptPass ())
        t_var.set(final_str)
        root.update_idletasks()
    return

#Login parameters in JIRA
def menuJIRA():
    #Ask for Users
    JIRAUser = simpledialog.askstring("JIRA USER", "Please introduce your username in JIRA", initialvalue= config['CONNECT'].get('Username'))
    print("New Username: " + JIRAUser)
    #Update User in Parameters
    config['CONNECT']['Username'] = JIRAUser
    with open(CurrentDir + '\ini\Parameters.ini', 'w') as configfile:
        config.write(configfile)
    #Ask for password
    JIRAPass1 = simpledialog.askstring("JIRA Password", "Please introduce your Password in JIRA", show='*')
    JIRAPass2 = simpledialog.askstring("Verification", "Confirm Password", show='*')
    if JIRAPass1 != JIRAPass2:
       messagebox.showinfo(title="ERROR", message="Passwords did not match, try again")
    else:
        Connect.CryptPass(JIRAPass1)
    messagebox.showinfo(title="SUCCESS", message="Login parameters saved")
    return

def Exit():
    sys.exit()
    return

def Version():
    # Update Screen Labels
    messagebox.showinfo(title="About", message=" Version : "+config['WINDOWS'].get('version'))
    final_str = '\n\n\n\n\n\n\n\n\n\n\n\n\n\n\nVersion: %s\nDeveloped by jose.camporro@vodafone.com' % \
                (config['WINDOWS'].get('version'))
    t_var.set(final_str)
    root.update()
    return

# ********************************
# Start Main Code Function
# ********************************

def StartPpt():
    # ********************************
    # Read Parameter again (if changes)
    # ********************************
    config = configparser.ConfigParser()
    config.sections()
    config.read(CurrentDir + '\ini\Parameters.ini')

    # CheckUsername
    if config['CONNECT'].get('Username') == "":
        menuJIRA()

    # ********************************
    # Function to Format Text
    # ********************************
    def formatTextStory(Texto, fsize):
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = Texto
        font = run.font
        font.name = config['INI'].get('Font')
        font.size = Pt(fsize)

    # ********************************
    # Function to duplicate a slide from template
    # ********************************
    def duplicate_slide(pres, index):
        template = pres.slides[index]
        try:
            blank_slide_layout = pres.slide_layouts[0]
        except:
            blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

        copied_slide = pres.slides.add_slide(blank_slide_layout)

        for shp in template.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for key, value in six.iteritems(template.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                copied_slide.part.rels.add_relationship(value.reltype,
                                                        value._target,
                                                        value.rId)
        return copied_slide

    # ********************************
    # Function to format Feature Text in the slide boxes
    # ********************************
    def formatText(Texto, fsize):
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = Texto
        font = run.font
        font.name = config['INI'].get('Font')
        font.size = Pt(fsize)

    # ********************************
    # Function to create QR Code
    # ********************************
    def InsertQRCode(issue):
        # Make QR Code
        img = qrcode.make(config['CONNECT'].get('JIRAserver') + "/browse/" + issue)

        # Adjust QR image
        basewidth = 85
        wpercent = (basewidth / float(img.size[0]))
        hsize = int((float(img.size[1]) * float(wpercent)))
        img = img.resize((basewidth, hsize), Image.ANTIALIAS)
        # Save QR File
        img.save("QR.jpg")

        # Add image in the slide
        slide.shapes.add_picture("QR.jpg", Inches(8.5), Inches(4.2))

    # ********************************
    # Function to calculate the total number of stories per feature
    # ********************************
    def StoryCalculator(Id):
        StoryQuery = config['INI'].get('StoryJQL') + " AND 'Epic Link'=" + Id

        storycount = 0
        try:
            # Execute Feature Query in JIRA
            Stories = jira.search_issues(StoryQuery, startAt=0, maxResults=False)
            storycount = len(Stories)
        except:
            storycount = 0
        return storycount


    # ********************************
    # Function to calculate the sum of story Points per feature
    # ********************************
    def StorySP(Id):
        StoryQuery = config['INI'].get('StoryJQL') + " and 'Epic Link'=" + Id

        storySP = 0

        # Execute Feature Query in JIRA
        Stories = jira.search_issues(StoryQuery, startAt=0, maxResults=False)

        for issues in Stories:
            if issues.fields.customfield_10002:
                storySP = storySP + issues.fields.customfield_10002
        return storySP

    # ********************************
    # Start Feature main programme
    # ********************************

    # Ini Counters
    TotFeat = 0
    TotStories = 0
    SUMStoryPoints = 0
    SUMStories = 0

    # Read Parameters to connect to JIRA
    jira_server = config['CONNECT'].get('JIRAserver')
    jira_user = config['CONNECT'].get('Username')
    jira_password = Connect.decryptPass ()

    # Update Screen Labels
    Status_var.set("Connecting JIRA")
    root.update()

    # Connect to JIRA
    jira_server = {'server': jira_server}
    jira = JIRA(options=jira_server, basic_auth=(jira_user, jira_password))

    # Prepare JIRA Query
    FeatQuery = config['INI'].get('FeatureJQL')

    # Execute Feature Query in JIRA
    Features = jira.search_issues(FeatQuery, startAt=0, maxResults=False)

    # Open PPTX
    pres = Presentation(CurrentDir + config['PPTX'].get('Template'))

    # Summaries of my features
    for issue in Features:
        # Count Features
        TotFeat = TotFeat + 1
        # Find Feature dependencies
        Dependencies = ""
        for link in issue.fields.issuelinks:
            if hasattr(link, "inwardIssue"):
                inwardIssue = link.inwardIssue
                Dependencies = inwardIssue.key + " " + Dependencies

        # Find Sprint
        try:
            for sprint in issue.fields.customfield_10004:
                sprint_name = re.findall(r"name=[^,]*", str(issue.fields.customfield_10004[0]))
                Sprint_text = str(sprint_name)
                Start_Sprint = Sprint_text.find("Sprint")
                Total_SprintLen = len(Sprint_text)
                # Get the right name for the sprint
                FeatSprint = Sprint_text[Start_Sprint:Total_SprintLen - 2]
        except:
            FeatSprint = "N/A"

        # Create Feature slide in PPTX
        title_slide_layout = pres.slide_layouts[0]
        slide = duplicate_slide(pres, 0)


        # Insert QR Code
        InsertQRCode(issue.key)

        # ********************************
        # Update Slide shapes
        # ********************************

        # Update the new Features slide
        for shape in slide.shapes:

            #Calculate the size of AC
            if len(str(issue.fields.customfield_11211)) > int(config['INI'].get('AC_Size')):
                Too_Big = True
            else:
                Too_Big = False

            if shape.name == 'FeatureName':
                formatText(issue.key + " " + issue.fields.summary, 14)
            if shape.name == 'Assignee':
                try:
                    formatText(issue.fields.assignee.displayName, 12)
                except:
                    formatText(" ", 12)
            #Change the title of the box
            if shape.name == "Text_Description":
                # If Description is True
                text_frame = shape.text_frame
                if config['INI'].get('Description') == "True" or Too_Big == True:

                    text_frame.paragraphs[0].runs[0].text = "Description"
                else:
                    text_frame.paragraphs[0].runs[0].text = "Acceptance Criteria"

            if shape.name == 'Description':
                # If Description is True
                if config['INI'].get('Description') == "True" or Too_Big == True:
                    Description = str(issue.fields.description).replace("\r", " ")
                else:
                    # Eliminate Line Return
                    Description = str(issue.fields.customfield_11211).replace("\r", " ")

                if len(Description) > 300:
                    Description = str(Description).replace("\n", " ")
                    # Copy Text
                formatText(Description, 11)
            if shape.name == 'Priority':
                formatText(str(issue.fields.priority), 12)
            if shape.name == 'Dependencies':
                formatText(Dependencies, 12)
            if shape.name == 'Sprint':
                formatText(FeatSprint, 12)
            # Sum of the Stories
            if shape.name == 'SUMStories':
                try:
                    SUMStories = str(StoryCalculator(issue.key))
                except:
                    SUMStories = 0
                formatText(SUMStories, 12)

            # Sum of the Story Points
            if shape.name == 'SUMStoryPoints':
                try:
                    SUMStoryPoints = str(StorySP(issue.key))
                except:
                    SUMStoryPoints = 0
                formatText(str(SUMStoryPoints), 12)
            # change color for team
            if shape.name == 'AreaBox':
                color = config['TEAMS'].get('Default.Color').split(",")
                for x in range(int(config['TEAMS'].get('Number.Team'))):
                    if config['TEAMS'].get(str(x) + '.Team.Name') in ascii(issue.fields.customfield_11018):
                        color = config['TEAMS'].get(str(x) + '.Team.Color').split(",")
                    shape.fill.fore_color.rgb = RGBColor(int(color[0]), int(color[1]), int(color[2]))
            # Update Team
            if shape.name == 'Area':
                teamfeature = " "
                for x in range(int(config['TEAMS'].get('Number.Team'))):
                    if config['TEAMS'].get(str(x) + '.Team.Name') in ascii(issue.fields.customfield_11018):
                        teamfeature = config['TEAMS'].get(str(x) + '.Team.Name')
                formatText(teamfeature, 12)

        # ********************************
        # Output for Features
        # ********************************

        #Output in console
        print('{}: {}: {}: {}'.format(issue.key, issue.fields.summary, FeatSprint, issue.fields.customfield_11018))

        # Update Screen Labels
        final_str = 'Feature Code: %s \nDescription: %s\nSprint: %s\nTeam: %s\n# Stories: %s\n# Story Points: %s' \
                    % (issue.key, issue.fields.summary, FeatSprint, issue.fields.customfield_11018, SUMStories, SUMStoryPoints)
        Status_var.set("Reading Features")
        t_var.set(final_str)
        root.update()

        # ********************************
        # Create Story slides in PPTX
        # ********************************
        if config['INI'].get('Stories') == "True":

            # ********************************
            # Read Stories in JIRA
            # ********************************
            StoryQuery = config['INI'].get('StoryJQL') + " AND 'Epic Link'=" + issue.key
            # Execute Feature Query in JIRA
            Stories = jira.search_issues(StoryQuery, startAt=0, maxResults=False)
            try:
                for issues in Stories:
                    # Add 1 to Stories Count
                    TotStories = TotStories + 1

                    # Create Story slide in PPTX
                    title_slide_layout = pres.slide_layouts[0]
                    slide = duplicate_slide(pres, 1)

                    # Find Story dependencies
                    Dependencies = ""
                    for link in issues.fields.issuelinks:
                        if hasattr(link, "inwardIssue"):
                            inwardIssue = link.inwardIssue
                            Dependencies = inwardIssue.key + " " + Dependencies

                    # Find Story Sprint
                    try:
                        for sprint in issues.fields.customfield_10004:

                              sprint_name = re.findall(r"name=[^,]*", str(issues.fields.customfield_10004[0]))
                              Sprint_text = str(sprint_name)
                              Start_Sprint = Sprint_text.find("Sprint")
                              Total_SprintLen = len(Sprint_text)
                              # Get the right name for the sprint
                              StorySprint = Sprint_text[Start_Sprint:Total_SprintLen - 2]
                    except:
                          StorySprint = ""

                    # Find Team of the Story
                    teamStory = " "
                    for x in range(int(config['TEAMS'].get('Number.Team'))):
                        if config['TEAMS'].get(str(x) + '.Team.Name') in ascii(issues.fields.customfield_11018):
                            teamStory = config['TEAMS'].get(str(x) + '.Team.Name')

                    # Insert QR Code
                    InsertQRCode(issues.key)

                    # ********************************
                    # Update the story slide
                    # ********************************
                    for shape in slide.shapes:
                        # Calculate the size of AC
                        if len(str(issues.fields.customfield_11211)) > int(config['INI'].get('AC_Size')):
                            Too_Big = True
                        else:
                            Too_Big = False

                        if shape.name == 'StoryName':
                            formatTextStory(issues.key + " " + issues.fields.summary, 14)
                        if shape.name == 'Assignee':
                            try:
                                formatTextStory(issues.fields.assignee.displayName, 12)
                            except:
                                formatTextStory(" ", 12)

                        # change the Title of the field
                        if shape.name == "Text_Description":
                            # If Description is True
                            text_frame = shape.text_frame
                            if config['INI'].get('Description') == "True" or Too_Big == True:

                                text_frame.paragraphs[0].runs[0].text = "Description"
                            else:
                                text_frame.paragraphs[0].runs[0].text = "Acceptance Criteria"

                        if shape.name == 'Description':
                            # If Description is True
                            if config['INI'].get('Description') == "True" or Too_Big == True:
                                Description = str(issues.fields.description).replace("\r", " ")
                            else:
                                # Eliminate Line Return
                                Description = str(issues.fields.customfield_11211).replace("\r", " ")

                            if len(Description) > 300:
                                Description = str(Description).replace("\n", " ")
                                # Copy Text
                            formatText(Description, 11)

                        if shape.name == 'Priority':
                            try:
                                formatTextStory(str(issues.fields.priority), 12)
                            except:
                                formatTextStory(" ", 12)
                        if shape.name == 'Dependencies':
                            try:
                                formatTextStory(Dependencies, 12)
                            except:
                                formatTextStory(" ", 12)
                        if shape.name == 'Sprint':
                            try:
                                formatTextStory(StorySprint, 12)
                            except:
                                formatTextStory(" ", 12)
                        # Story Points
                        if shape.name == 'StoryPoints':
                            try:
                                StoryPoints = str(issues.fields.customfield_10002)
                            except:
                                StoryPoints = 0
                            formatTextStory(StoryPoints, 12)
                        # Link to Feature
                        if shape.name == 'Parent':
                            try:
                                formatTextStory(issue.key + " " + issue.fields.summary, 12)
                            except:
                                formatTextStory(" ", 12)
                        # Same Team as the Feature
                        if shape.name == 'Area':
                            formatText(teamStory, 12)
                        # change color for team
                        if shape.name == 'AreaBox':
                            try:
                                color = config['TEAMS'].get('Default.Color').split(",")
                                for x in range(int(config['TEAMS'].get('Number.Team'))):
                                    if config['TEAMS'].get(str(x) + '.Team.Name') in ascii(
                                            issues.fields.customfield_11018):
                                        color = config['TEAMS'].get(str(x) + '.Team.Color').split(",")
                            except:
                                color = config['TEAMS'].get('0.Team.Color').split(",")
                            shape.fill.fore_color.rgb = RGBColor(int(color[0]), int(color[1]), int(color[2]))
                # ********************************
                # Output for Stories
                # ********************************

                #Output in console
                print('    {}: {}: {}: {}'.format(issues.key, issues.fields.summary, StorySprint, issues.fields.customfield_11018))

                # Update Screen Labels
                final_str = 'Story Code: %s \nDescription: %s\nSprint: %s\nStory Points: %s' \
                            % (issues.key, issues.fields.summary,  StorySprint, StoryPoints)
                Status_var.set("Reading Stories")
                t_var.set(final_str)
                root.update()
            except:
                print("        No Story")

    # ********************************
    # Close PowerPoints
    # ********************************

    # Update Screen Labels
    Status_var.set("writing PowerPoint")
    root.update_idletasks()
    try:
        pres.save(CurrentDir + config['PPTX'].get('Output'))
    except Exception as e:
        t_var.set(str(e))
        Status_var.set("Error Writing the Ppt file")
    print("PowerPoint created")

    # ********************************
    # End Process
    # ********************************

    # Update Screen Labels
    Status_var.set("PowerPoint Created")
    final_str = 'File created: %s \n' \
                'Total Features: %s \nTotal Stories: %s' % \
                (config['PPTX'].get('Output'), TotFeat, TotStories)
    t_var.set(final_str)
    os.remove("QR.jpg")
    root.update()
    #Make Visible Powerpoint button
    Ppt_button.place(relx=0.5, rely=0.5, anchor=CENTER)
    Underframe.place(relx=0.5, rely=0.9, relwidth=0.23, relheight=0.07, anchor='n')

# ********************************
# Start Main Code
# ********************************

if __name__ == '__main__':

    # Initiate Parameters
    CurrentDir = os.getcwd()

    # Create Class object
    Connect = Connect()
    config = Connect.GetConfig()
    ConfReports = ConfReports()

    #Parameters Ini
    config = configparser.ConfigParser()
    config.sections()
    config.read(CurrentDir + '\ini\Parameters.ini')

    #Initiate Windows
    root = Tk()
    root.title(config['WINDOWS'].get('Title'))

    #Create Menu
    menu = Menu(root)
    root.config(menu=menu)
    ConfigMenu = Menu(menu)

    #Initial Menu Option
    menu.add_cascade(label="Config", menu=ConfigMenu)
    ConfigMenu.add_command(label="Settings", command=Parameters)
    ConfigMenu.add_command(label="Login JIRA", command=menuJIRA)
    ConfigMenu.add_command(label="Check Connection", command=CheckJIRA)
    ConfigMenu.add_separator()
    ConfigMenu.add_command(label="Exit", command=Exit)

    ConfluenceMenu = Menu(menu)
    menu.add_cascade(label="Confluence", menu= ConfluenceMenu )
    ConfluenceMenu.add_command(label="PI Programme Board ", command=PIBoard)
    ConfluenceMenu.add_command(label="Team Report", command=TeamBoard)
    ConfluenceMenu.add_command(label="Assignee Report", command=AssigneeBoard)
    ConfluenceMenu.add_command(label="Pending Features Report", command=PendingFeatures)


    #About Menu option
    AboutMenu = Menu(menu)
    menu.add_cascade(label="About", menu= AboutMenu )
    AboutMenu.add_command(label="Version", command=Version)


    canvas = Canvas(root, height=config['WINDOWS'].get('height'), width = config['WINDOWS'].get('width'))
    canvas.pack()


    background_image = PhotoImage(file=CurrentDir + '\image\one.png')
    background_label = Label(root , image = background_image)
    background_label.place(relwidth=1, relheight=1)

    frame = Frame(root, bg=config['WINDOWS'].get('framecolor'), bd=5)
    frame.place(relx=0.5, rely=0.1, relwidth=0.75, relheight=0.1, anchor='n')

    # Init text variable for lower Frame
    Status_var = StringVar()
    Status_var.set("Click on Start Button")

    status = Label(frame,  textvariable=Status_var, font=40)
    status.place(relwidth=0.65, relheight=1)

    B = Button(frame, text= "Start", font=40, command=StartPpt)
    B.place(relx=0.7, relwidth=0.3, relheight=1)

    lower_frame = Frame(root, bg=config['WINDOWS'].get('framecolor'), bd=10)
    lower_frame.place(relx=0.5, rely=0.25, relwidth=0.75, relheight=0.6, anchor='n')

    # Init text variable for lower Frame,
    t_var = StringVar()

    theLabel = Label(lower_frame, font= 20, textvariable=t_var, anchor='nw', justify='left', wraplength=500)
    theLabel.place(relwidth=1, relheight=1)

    #Button Hidden
    Underframe = Frame(root, bg=config['WINDOWS'].get('framecolor'), bd=5)
    Ppt_button = Button(Underframe, text='Open PowerPoint',font=20, command=Openppt)

    root.mainloop()







